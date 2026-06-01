"""Генерирует тестовую PPTX-презентацию с профилем сотрудника."""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "samples" / "sample_employee.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, lines):
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    left = Inches(1); top = Inches(2); width = Inches(11); height = Inches(4)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
    return slide

# Slide 1 — Шапка
add_title_slide("Профиль сотрудника", [
    "ФИО: Петров Пётр Петрович",
    "Должность: Старший мастер участка",
    "Подразделение: Цех молочной продукции №3",
    "Возраст: 42",
    "Пол: мужской",
    "Стаж: 12",
    "Образование: высшее техническое",
])

# Slide 2 — 16PF
add_title_slide("16PF Кеттелла (стэны 1-10)", [
    "A: 7   B: 6   C: 5   E: 8",
    "F: 5   G: 9   H: 6   I: 5",
    "L: 4   M: 5   N: 6   O: 3",
    "Q1: 5  Q2: 7  Q3: 7  Q4: 4",
])

# Slide 3 — Big Five
add_title_slide("Big Five (баллы 0-100)", [
    "Открытость: 65",
    "Добросовестность: 78",
    "Экстраверсия: 55",
    "Доброжелательность: 70",
    "Нейротизм: 35",
])

# Slide 4 — MMPI
add_title_slide("MMPI / СМИЛ (T-баллы)", [
    "Валидность: L: 52  F: 48  K: 56",
    "Клинические:",
    "Hs: 55  D: 62  Hy: 50  Pd: 48",
    "Mf: 50  Pa: 55  Pt: 60  Sc: 45",
    "Ma: 50  Si: 55",
])

# Slide 5 — DISC
add_title_slide("DISC", [
    "D: 70  I: 50  S: 60  C: 75",
])

# Slide 6 — HOLLAND
add_title_slide("HOLLAND RIASEC", [
    "R: 30  I: 70  A: 45  S: 65  E: 80  C: 50",
    "Код: ESC",
])

# Slide 7 — MBTI
add_title_slide("MBTI", [
    "Тип: INTJ",
])

# Slide 8 — Амтхауэр
add_title_slide("Амтхауэр", [
    "IQ: 115",
])

prs.save(OUT)
print(f"✅ Сохранено: {OUT}")
print(f"   Размер: {OUT.stat().st_size} байт, {len(prs.slides)} слайдов")
