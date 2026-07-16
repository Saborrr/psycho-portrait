"""Структурное чтение карточек ЭФКО из PPTX.

В отличие от общего regex-парсера этот модуль читает таблицы и данные диаграмм
напрямую из Office Open XML. Это сохраняет связи «показатель - значение» и не
зависит от порядка shape-объектов внутри слайда.
"""
from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path
from typing import Iterable, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from .models import MMPI, ParsedProfile
from .methods.efko import (
    AchievementModel,
    Aktivnost,
    ConsciousUnconsciousScore,
    EchvEchm,
    EFKOSet,
    Empaty,
    EmployerParadigm,
    HarmonyDecision,
    IntellektEfk,
    LifeGamble,
    PersonnelTypes,
    SafetyAttitude,
    WorkDiscomfort,
)


NAME_RE = re.compile(r"^([А-ЯЁ][А-ЯЁ-]+(?:\s+[А-ЯЁ][А-ЯЁ-]+){2,3})$")
ORG_RE = re.compile(r"\b(?:ООО|АО|ОАО|ЗАО|ПАО|ФИЛИАЛ)\b", re.IGNORECASE)


def _clean(value: object) -> str:
    return re.sub(r"\s+", " ", str(value or "").replace("\xa0", " ")).strip()


def _key(value: object) -> str:
    value = _clean(value).lower().replace("ё", "е")
    return re.sub(r"[^a-zа-я0-9]+", " ", value).strip()


def _number(value: object) -> Optional[float]:
    match = re.search(r"[-+]?\d+(?:[.,]\d+)?", _clean(value))
    if not match:
        return None
    try:
        return float(match.group(0).replace(",", "."))
    except ValueError:
        return None


def _percent(value: object) -> Optional[float]:
    number = _number(value)
    if number is None:
        return None
    # Данные диаграмм Office часто хранят проценты как 0.67, а таблицы как 67.
    if 0 <= number <= 1:
        number *= 100
    return round(number, 2)


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()


def _rows(shape) -> list[list[str]]:
    return [[_clean(cell.text) for cell in row.cells] for row in shape.table.rows]


def _all_tables(prs: Presentation) -> Iterable[list[list[str]]]:
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = _rows(shape)
                if rows:
                    yield rows


def _set_if_number(obj, field: str, value: object) -> bool:
    number = _number(value)
    if number is None:
        return False
    setattr(obj, field, number)
    return True


def _parse_employee_headers(prs: Presentation, profile: ParsedProfile, source_filename: str | None) -> None:
    top_texts: list[tuple[int, int, str]] = []
    # ФИО и должность находятся на первом слайде. Заголовки следующих слайдов
    # (например «Азарт к жизни») нельзя принимать за должность сотрудника.
    first_slide = prs.slides[0]
    for shape in first_slide.shapes:
        text = _shape_text(shape)
        if text and shape.top < Inches(1.0):
            top_texts.append((int(shape.top), int(shape.left), text))

    employee_header: str | None = None
    for _, _, text in sorted(top_texts):
        first_line = _clean(text.splitlines()[0])
        match = NAME_RE.match(first_line)
        if match and re.search(r"\b\d{1,2}\s*(?:лет|год|года)\b", text, re.IGNORECASE):
            profile.employee.full_name = match.group(1).title()
            employee_header = text
            if age := re.search(r"\b(\d{1,2})\s*(?:лет|год|года)\b", text, re.IGNORECASE):
                profile.employee.age = int(age.group(1))
            break

    if not profile.employee.full_name and source_filename:
        stem = Path(source_filename).stem.strip()
        if NAME_RE.match(stem):
            profile.employee.full_name = stem.title()
            profile.notes.append("ФИО определено из имени файла.")

    # Левый верхний заголовок содержит должность и организацию.
    for _, left, text in sorted(top_texts):
        if left > Inches(6.0) or text == employee_header:
            continue
        compact = _clean(text)
        if len(compact) < 8 or compact.isdigit():
            continue
        org_match = ORG_RE.search(compact)
        if org_match:
            profile.employee.position = compact[:org_match.start()].strip(" ,")
            profile.employee.department = compact[org_match.start():].strip()
        elif not profile.employee.position:
            profile.employee.position = compact
        break


def _parse_employee_table(rows: list[list[str]], profile: ParsedProfile) -> bool:
    labels = {_key(row[0]) for row in rows if row}
    if not any(label.startswith("возраст место проживания") for label in labels):
        return False
    for row in rows:
        if len(row) < 2:
            continue
        label, value = _key(row[0]), row[1].strip()
        if label.startswith("возраст место проживания"):
            if age := re.search(r"\b(\d{1,2})\s*(?:лет|год|года)\b", value, re.IGNORECASE):
                profile.employee.age = int(age.group(1))
            profile.employee.extra["location"] = value
        elif label.startswith("семейное положение"):
            profile.employee.extra["family_status"] = value
            low = value.lower()
            if any(marker in low for marker in ("замуж", "разведена", "вдова")):
                profile.employee.gender = "женский"
            elif any(marker in low for marker in ("женат", "не женат", "холост", "разведен", "вдовец")):
                profile.employee.gender = "мужской"
        elif label.startswith("образование"):
            profile.employee.education = value or None
        elif label.startswith("опыт работы"):
            profile.employee.extra["employment_history"] = value
        elif label.startswith("готовность к командировкам"):
            if value and value not in {"-", "—"}:
                profile.employee.extra["mobility"] = value
    return True


MMPI_LABELS = {
    "l": "L", "l лжи": "L", "f": "F", "k": "K", "k коррекции": "K",
    "1 ипохондрия": "Hs", "2 депрессия": "D", "3 истероидность": "Hy",
    "4 психопатия": "Pd", "5 феминизированность": "Mf", "6 паранойяльность": "Pa",
    "7 психастения": "Pt", "8 шизоидность": "Sc", "9 гипомания": "Ma",
    "9 маниакальность": "Ma", "0 соц интроверсия": "Si",
    "0 социальная интроверсия": "Si", "джентльмен": "gentleman",
    "здравомыслие": "sanity",
}


def _parse_mmpi_table(rows: list[list[str]], profile: ParsedProfile) -> bool:
    if not rows or "пси характеристика" not in _key(rows[0][0]):
        return False
    data = MMPI()
    found = 0
    for row in rows[1:]:
        if len(row) < 2:
            continue
        field = MMPI_LABELS.get(_key(row[0]))
        if field and _set_if_number(data, field, row[1]):
            found += 1

    # Правила достоверности из лекции НОЦ «Бирюч», актуализация 30.01.2025.
    validity_values = [v for v in (data.L, data.F, data.K) if v is not None]
    if validity_values and any(v > 80 for v in validity_values):
        data.validity = "invalid"
    elif validity_values and any(v >= 70 for v in validity_values):
        data.validity = "questionable"
    elif len(validity_values) == 3:
        data.validity = "valid"

    peaks = []
    for code, field in [("1", "Hs"), ("2", "D"), ("3", "Hy"), ("4", "Pd"),
                        ("5", "Mf"), ("6", "Pa"), ("7", "Pt"), ("8", "Sc"),
                        ("9", "Ma"), ("0", "Si")]:
        value = getattr(data, field)
        if value is not None and value >= 70:
            peaks.append((value, code))
    peaks.sort(reverse=True)
    if peaks:
        data.code = "-".join(code for _, code in peaks[:3])
        data.profile_type = {1: "пиковый", 2: "двухпиковый"}.get(len(peaks), "трёхпиковый")

    profile.methods.mmpi = data
    if found < 15:
        profile.notes.append(f"СМИЛ: структурно распознано {found}/15 числовых показателей.")
    return True


def _parse_efko_table(rows: list[list[str]], efko: EFKOSet) -> bool:
    if not rows:
        return False
    header = _key(rows[0][0])
    if header == "интеллект":
        data = IntellektEfk()
        mapping = {"логический интеллект": "logicheskiy", "образный интеллект": "obrazny", "лексика": "leksika"}
        for row in rows[1:]:
            if len(row) >= 2 and (field := mapping.get(_key(row[0]))):
                setattr(data, field, _percent(row[1]))
        efko.intellekt = data
        return True
    if header == "активность":
        data = Aktivnost()
        mapping = {"физическая": "fizicheskaya", "интеллектуальная": "intellektualnaya", "коммуникационная": "kommunikacionnaya"}
        for row in rows[1:]:
            if len(row) >= 2 and (field := mapping.get(_key(row[0]))):
                setattr(data, field, _percent(row[1]))
        efko.aktivnost = data
        return True
    if header == "эмпатия":
        data = Empaty()
        for row in rows[1:]:
            if _key(row[0]).startswith("2 рода"):
                if len(row) >= 2: data.conscious = _percent(row[1])
                if len(row) >= 3: data.unconscious = _percent(row[2])
        efko.empaty = data
        return True
    if header.startswith("постмодерн"):
        if len(rows[0]) >= 2:
            efko.postmodern = _percent(rows[0][1])
        return True
    if header.startswith("отношение к дискомфорту в работе"):
        data = WorkDiscomfort()
        mapping = {"втр": "vtr", "сэдо": "sedo", "седо": "sedo", "опрятность": "neatness", "норматив командировок": "business_trips"}
        for row in rows[1:]:
            if len(row) >= 2 and (field := mapping.get(_key(row[0]))):
                setattr(data, field, _percent(row[1]))
        efko.work_discomfort = data
        return True
    if header.startswith("модель по отношению к правилам безопасности"):
        data = SafetyAttitude()
        mapping = {"обезбашенный герой": "reckless_hero", "безынициативный исполнитель": "passive_executor", "защитник": "protector"}
        for row in rows[1:]:
            if len(row) >= 2 and (field := mapping.get(_key(row[0]))):
                setattr(data, field, _percent(row[1]))
        efko.safety_attitude = data
        return True
    return False


def _chart_data(shape) -> tuple[list[str], list[tuple[str, list[Optional[float]]]]]:
    chart = shape.chart
    try:
        categories = [_clean(c.label) for c in chart.plots[0].categories]
    except Exception:
        return [], []
    series = []
    for item in chart.series:
        series.append((_clean(item.name), [_percent(v) for v in item.values]))
    return categories, series


def _first_series_values(series: list[tuple[str, list[Optional[float]]]]) -> list[Optional[float]]:
    return series[0][1] if series else []


def _value_at(values: list[Optional[float]], index: int) -> Optional[float]:
    return values[index] if index < len(values) else None


def _parse_chart(categories: list[str], series: list[tuple[str, list[Optional[float]]]], efko: EFKOSet) -> bool:
    keys = [_key(c) for c in categories]
    if keys == ["эчв", "эчм соц отношений"]:
        values = _first_series_values(series)
        efko.echv_echm = EchvEchm(echv=_value_at(values, 0), echm=_value_at(values, 1))
        return True
    if keys == ["ничего личного", "партнерство", "солидарность"]:
        by_name = {_key(name): values for name, values in series}
        conscious = by_name.get("сознательная оценка", [])
        unconscious = by_name.get("бессознательная оценка", [])
        efko.employer_paradigm = EmployerParadigm(
            nothing_personal=ConsciousUnconsciousScore(conscious=_value_at(conscious, 0), unconscious=_value_at(unconscious, 0)),
            partnership=ConsciousUnconsciousScore(conscious=_value_at(conscious, 1), unconscious=_value_at(unconscious, 1)),
            solidarity=ConsciousUnconsciousScore(conscious=_value_at(conscious, 2), unconscious=_value_at(unconscious, 2)),
        )
        return True
    if keys == ["созидательный азарт", "праздный азарт", "пассивность"]:
        values = _first_series_values(series)
        efko.life_gamble = LifeGamble(creative=_value_at(values, 0), idle=_value_at(values, 1), passivity=_value_at(values, 2))
        return True
    if keys == ["гармония дискомфорта", "решительность", "гармония радости", "осторожность"]:
        values = _first_series_values(series)
        efko.harmony_decision = HarmonyDecision(
            discomfort_harmony=_value_at(values, 0), decisiveness=_value_at(values, 1),
            joy_harmony=_value_at(values, 2), caution=_value_at(values, 3),
        )
        return True
    if len(keys) == 3 and keys[0].startswith("рационально предприимчивый"):
        values = _first_series_values(series)
        efko.achievement_model = AchievementModel(
            rational_enterprising=_value_at(values, 0), enabling_entrepreneurs=_value_at(values, 1),
            public_dismissal_significance=_value_at(values, 2),
        )
        return True
    if keys == ["труженик", "эволюционное развитие", "лидершип", "поединщик"]:
        values = _first_series_values(series)
        efko.personnel_types = PersonnelTypes(
            worker=_value_at(values, 0), evolutionary_development=_value_at(values, 1),
            leadership=_value_at(values, 2), duelist=_value_at(values, 3),
        )
        return True
    return False


def enrich_profile_from_pptx(path: str, profile: ParsedProfile, source_filename: str | None = None) -> ParsedProfile:
    prs = Presentation(path)
    profile.source_filename = source_filename or Path(path).name
    _parse_employee_headers(prs, profile, source_filename or Path(path).name)

    efko = profile.methods.efko or EFKOSet()
    # Удаляем значения, которые старый свободный regex мог захватить из соседнего блока.
    efko.paradigma_k_rabotodatelyu = None
    efko.gotovnost_k_komandirovkam = None
    structured_tables = 0
    for rows in _all_tables(prs):
        if _parse_employee_table(rows, profile):
            structured_tables += 1
        elif _parse_mmpi_table(rows, profile):
            structured_tables += 1
        elif _parse_efko_table(rows, efko):
            structured_tables += 1

    structured_charts = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                categories, series = _chart_data(shape)
                if _parse_chart(categories, series, efko):
                    structured_charts += 1

    if structured_tables or structured_charts:
        profile.methods.efko = efko
        profile.notes = [
            n for n in profile.notes
            if not n.startswith("ЭФКО:") and not n.startswith("MMPI/СМИЛ:")
        ]
        profile.notes.append(
            f"Структурный PPTX-парсер: обработано таблиц {structured_tables}, диаграмм {structured_charts}."
        )
    return profile


def extract_employee_photo(path: str) -> tuple[bytes, str] | None:
    """Возвращает фото сотрудника и расширение, не включая его в JSON профиля."""
    prs = Presentation(path)
    candidates = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            score = 0
            if "фото сотрудника" in shape.name.lower():
                score += 100
            score += int(shape.width * shape.height / 1_000_000_000_000)
            candidates.append((score, shape))
    if not candidates:
        return None
    _, shape = max(candidates, key=lambda item: item[0])
    blob = shape.image.blob
    try:
        from PIL import Image
        with Image.open(BytesIO(blob)) as image:
            if image.width * image.height > 25_000_000:
                return None
            image.verify()
    except Exception:
        return None
    return blob, shape.image.ext
