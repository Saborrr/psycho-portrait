"""
Парсер PPTX-презентаций с результатами психологических тестов.

Стратегия:
1. Извлечь весь текст + таблицы со слайдов.
2. Эвристически определить шапку (ФИО, должность, возраст, подразделение).
3. Найти упоминания методик (16PF, Big Five, MMPI, DISC, HOLLAND, MBTI, Амтхауэр).
4. Для каждой методики — пройтись по паттернам «фактор → число».

Шаблоны презентации см. в docs/presentation_template.md.
"""
from __future__ import annotations
import os
import re
from typing import Optional
from pptx import Presentation

from .models import (
    EmployeeInfo, MethodScores, ParsedProfile,
    Cattell16PF, BigFive, MMPI, DISC, Holland, MBTI, Amthauer,
)


# === Утилиты ===

def _iter_shapes(slide):
    """Рекурсивно обойти все shape-ы на слайде (включая группы и таблицы)."""
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GROUP
            yield from _iter_shapes(shape)
        else:
            yield shape


def _shape_text(shape) -> str:
    """Текст из shape-а (включая таблицы)."""
    parts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            # Runs являются фрагментами форматирования одного абзаца.
            parts.append("".join(run.text for run in para.runs) or para.text)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(p for p in parts if p)


def _chart_text(shape) -> str:
    """Аудируемое текстовое представление встроенной диаграммы."""
    if not getattr(shape, "has_chart", False):
        return ""
    try:
        categories = [str(c.label).strip() for c in shape.chart.plots[0].categories]
        lines = ["ДИАГРАММА: " + " | ".join(categories)]
        for series in shape.chart.series:
            values = []
            for value in series.values:
                number = float(value)
                if 0 <= number <= 1:
                    number *= 100
                values.append(f"{number:g}")
            lines.append(f"{str(series.name or 'значения').strip()}: " + " | ".join(values))
        return "\n".join(lines)
    except Exception:
        return ""


def extract_all_text(path: str) -> tuple[str, list[str], int]:
    """Возвращает (объединённый_текст, текст_по_слайдам, количество_слайдов)."""
    prs = Presentation(path)
    max_slides = max(1, int(os.getenv("MAX_SLIDES", "50")))
    if len(prs.slides) > max_slides:
        raise ValueError(f"PPTX содержит больше {max_slides} слайдов")
    slides_text = []
    for slide in prs.slides:
        chunks = []
        for shape in _iter_shapes(slide):
            t = _shape_text(shape)
            if t.strip():
                chunks.append(t.strip())
            chart = _chart_text(shape)
            if chart:
                chunks.append(chart)
        slides_text.append("\n".join(chunks))
    full = "\n\n--- SLIDE BREAK ---\n\n".join(slides_text)
    return full, slides_text, len(prs.slides)


# === Шапка: ФИО, должность и т.п. ===

NAME_RE = re.compile(r"(?:ФИО|Сотрудник|Респондент|Имя|Фамилия)[ \t:]+([А-ЯЁ][а-яё\-]+(?:\s+[А-ЯЁ][а-яё\-]+){1,3})(?=\s*(?:\n|$|,|;))", re.IGNORECASE)
POS_RE = re.compile(r"(?:Должность|Позиция|Position)[ \t:]+([^\n\r,;]+)", re.IGNORECASE)
DEPT_RE = re.compile(r"(?:Подразделение|Отдел|Департамент|Цех|Участок|Филиал)[ \t:]+([^\n\r,;]+)", re.IGNORECASE)
AGE_RE = re.compile(r"(?:Возраст|Дата рождения)[ \t:]+(?:от\s*)?(\d{1,2}(?:\s*лет)?|\d{4}-\d{2}-\d{2}|\d{1,2}\.\d{1,2}\.\d{4})", re.IGNORECASE)
GENDER_RE = re.compile(r"(?:Пол|Гендер)[ \t:]+(мужской|женский|М|Ж|male|female)", re.IGNORECASE)
EDU_RE = re.compile(r"(?:Образование|Уровень образования)[ \t:]+([^\n\r]+)", re.IGNORECASE)
TENURE_RE = re.compile(r"(?:Стаж(?:\s+работы)?|Опыт(?:\s+работы)?)[ \t:]+(\d+(?:[.,]\d+)?)", re.IGNORECASE)


def parse_employee_info(text: str) -> EmployeeInfo:
    info = EmployeeInfo()
    if m := NAME_RE.search(text):
        info.full_name = m.group(1).strip().rstrip(",;")
    if m := POS_RE.search(text):
        info.position = m.group(1).strip()
    if m := DEPT_RE.search(text):
        info.department = m.group(1).strip()
    if m := AGE_RE.search(text):
        age_raw = m.group(1)
        digits = re.search(r"\d{1,2}", age_raw)
        if digits:
            info.age = int(digits.group(0))
    if m := GENDER_RE.search(text):
        g = m.group(1).lower()
        info.gender = "мужской" if g.startswith(("м", "m")) else "женский"
    if m := EDU_RE.search(text):
        info.education = m.group(1).strip()
    if m := TENURE_RE.search(text):
        info.tenure_years = float(m.group(1).replace(",", ".").split()[0])
    return info


# === Признаки методик ===

METHOD_HINTS = {
    "cattell_16pf": [
        r"\b16\s*PF\b", r"\b16PF\b", r"Кеттелл", r"Кэттелл", r"Cattell",
    ],
    "big_five": [
        r"Big\s*Five", r"Большая\s*пят[её]рка", r"NEO[\s\-]PI", r"BFI\b", r"Big-5",
    ],
    "mmpi": [
        r"\bMMPI\b", r"\bСМИЛ\b", r"Миннесотский", r"Миннесо?т",
        r"ПСИ\s*ХАРАКТЕРИСТИКА",
    ],
    "disc": [
        r"\bDISC\b", r"Диск[\- ]тест", r"Дисковая",
    ],
    "holland": [
        r"\bHOLLAND\b", r"\bХолланд\b", r"RIASEC", r"Холл?анд",
    ],
    "mbti": [
        r"\bMBTI\b", r"Майерс[ \-]?Бриггс", r"Майерс\-Бригс", r"ТИМ\b",
    ],
    "amthauer": [
        r"Амтхауэр", r"Амтауэр", r"структур[аы] интеллект[ау]",
    ],
}


def detect_method_blocks(text: str) -> dict[str, str]:
    """Для каждой методики — вырезать кусок текста, где она упоминается."""
    blocks = {k: "" for k in METHOD_HINTS}
    # Делим текст на абзацы/строки
    lines = text.splitlines()
    for i, line in enumerate(lines):
        for method, patterns in METHOD_HINTS.items():
            for pat in patterns:
                if re.search(pat, line, re.IGNORECASE):
                    # Длинные таблицы СМИЛ занимают более 30 строк после
                    # преобразования ячеек в текст.
                    start = max(0, i - 1)
                    end = min(len(lines), i + 45)
                    blocks[method] += "\n" + "\n".join(lines[start:end])
                    break
    return blocks


# === Парсеры по методикам ===

def _num(s: str) -> Optional[float]:
    """Достать первое число из строки (в т.ч. '7,5' → 7.5)."""
    m = re.search(r"[-+]?\d+(?:[.,]\d+)?", s.replace(",", "."))
    if m:
        try:
            return float(m.group(0))
        except ValueError:
            return None
    return None


def parse_cattell(block: str) -> tuple[Cattell16PF, list[str]]:
    """16PF — ищем паттерны вида 'A: 7', 'Фактор A — 7 стэнов', 'A  7', 'A=7'."""
    data = Cattell16PF()
    notes = []
    keys = list(Cattell16PF.model_fields.keys())  # A, B, C, E, F, G, H, I, L, M, N, O, Q1, Q2, Q3, Q4
    found = 0
    for k in keys:
        # Регулярка: \bK\s*[:=—\-]?\s*число
        pat = rf"\b{k}\b\s*[:=—\-]?\s*(\d{{1,2}}(?:[.,]\d+)?)"
        m = re.search(pat, block, re.IGNORECASE)
        if m:
            val = _num(m.group(1))
            if val is not None and 0 <= val <= 10:
                setattr(data, k, val)
                found += 1
    if found == 0:
        notes.append("16PF: не удалось распознать ни одного фактора. Проверь шаблон.")
    elif found < 8:
        notes.append(f"16PF: распознано только {found}/16 факторов. Возможно, формат нестандартный.")
    return data, notes


def parse_big_five(block: str) -> tuple[BigFive, list[str]]:
    data = BigFive()
    notes = []
    # Синонимы факторов
    field_synonyms = {
        "openness":         [r"открытость", r"openness", r"\bO\b", r"открытость опыту"],
        "conscientiousness":[r"добросовестность", r"conscientiousness", r"ответственность", r"\bC\b"],
        "extraversion":     [r"экстраверсия", r"extraversion", r"общительность", r"\bE\b"],
        "agreeableness":    [r"доброжелательность", r"уживчивость", r"согласие", r"agreeableness", r"\bA\b"],
        "neuroticism":      [r"нейротизм", r"neuroticism", r"эмоц\.\s*нестабильность", r"\bN\b"],
    }
    found = 0
    for field, syns in field_synonyms.items():
        for syn in syns:
            # 'Открытость: 65' / 'Openness = 65' / 'O = 65'
            pat = rf"(?:{syn})\s*[:=—\-]?\s*(\d{{1,3}}(?:[.,]\d+)?)"
            m = re.search(pat, block, re.IGNORECASE)
            if m:
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 100:
                    setattr(data, field, val)
                    found += 1
                    break
    if found == 0:
        notes.append("Big Five: ничего не распознано. Проверь формат.")
    return data, notes


def parse_mmpi(block: str) -> tuple[MMPI, list[str]]:
    """Парсер СМИЛ / MMPI / ММИЛ.
    Распознаёт:
      - Валидность L/F/K
      - Базисные шкалы Hs/D/Hy/Pd/Mf/Pa/Pt/Sc/Ma/Si
      - Код профиля (напр. '2-4-7', 'Профиль: 4-9', 'Код: 2-4-7')
      - Тип профиля ('линейный', 'пиковый', 'смешанный')
    """
    data = MMPI()
    notes = []
    # Синонимы для каждой шкалы — название + лексема
    scale_synonyms = {
        "L":  [r"\bложь\b", r"\bлжи\b", r"\bLie\b", r"\bL-?шкала\b"],
        "F":  [r"\bдостоверность\b", r"\bF-?шкала\b", r"\bаггравация\b", r"\bFrequency\b"],
        "K":  [r"\bкоррекция\b", r"\bK-?шкала\b", r"\bCorrection\b", r"\bзащитн\w+\s+установк"],
        "Hs": [r"\bипохондрия\b", r"\bHs\b", r"\bсверхконтрол"],
        "D":  [r"\bпессимистичност", r"\bдепрессия\b", r"\bD\b", r"\bшкала\s*2\b"],
        "Hy": [r"\bэмоциональн\w+\s+лабильност", r"\bистерия\b", r"\bконверсия\b", r"\bHy\b"],
        "Pd": [r"\bимпульсивност", r"\bасоциальн\w+\s+психопатия", r"\bPd\b"],
        "Mf": [r"\bмаскулинност", r"\bфемининност", r"\bMf\b"],
        "Pa": [r"\bригидност", r"\bпаранойя\b", r"\bPa\b"],
        "Pt": [r"\bтревожност", r"\bпсихастения\b", r"\bPt\b"],
        "Sc": [r"\bиндивидуалистичност", r"\bшизофрения\b", r"\bSc\b"],
        "Ma": [r"\bоптимизм\w+\s+активност", r"\bгипомания\b", r"\bMa\b"],
        "Si": [r"\bсоциальн\w+\s+интроверс", r"\bSi\b"],
    }
    found = 0
    for k, syns in scale_synonyms.items():
        for syn in syns:
            pat = rf"(?:{syn})[ \t]*[:=—\-][ \t]*(\d{{1,3}}(?:[.,]\d+)?)"
            m = re.search(pat, block, re.IGNORECASE)
            if m:
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 120:
                    setattr(data, k, val)
                    found += 1
                    break

    # Если для базисных шкал не сработали длинные синонимы — пробуем одиночные коды
    # Но осторожно: D/Pt/Sc — пересекаются с другими. Ищем только в формате 'D: 65' и т.п.
    code_only = {
        "L":  r"\bL\b",  "F":  r"\bF\b",  "K":  r"\bK\b",
        "Hs": r"\bHs\b", "D":  r"\bD\b",  "Hy": r"\bHy\b",
        "Pd": r"\bPd\b", "Mf": r"\bMf\b", "Pa": r"\bPa\b",
        "Pt": r"\bPt\b", "Sc": r"\bSc\b", "Ma": r"\bMa\b", "Si": r"\bSi\b",
    }
    for k, code_pat in code_only.items():
        if getattr(data, k) is None:
            pat = rf"(?:{code_pat})[ \t]*[:=—\-][ \t]*(\d{{1,3}}(?:[.,]\d+)?)"
            m = re.search(pat, block)
            if m:
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 120:
                    setattr(data, k, val)
                    found += 1

    # Код профиля
    code_m = re.search(r"\b(?:код|профиль)[\s:]+([1-9](?:-[1-9]){0,2})\b", block, re.IGNORECASE)
    if code_m:
        data.code = code_m.group(1)

    # Тип профиля
    type_m = re.search(r"\bтип[\s:]+профиля[\s:]+(линейный|пиковый|смешанный|двухпиковый|трёхпиковый|конверсионный|невротический|психотический|амбивалентный)\b", block, re.IGNORECASE)
    if type_m:
        data.profile_type = type_m.group(1).lower()

    # Корпоративные оценочные шкалы (ЭФКО и т.п.)
    # Джентльмен — социально-корректное поведение, этикет, следование конвенциональным нормам
    gentleman_m = re.search(r"\b(?:джентльмен|gentleman)\b[ \t]*[:=—\-][ \t]*(\d{1,3}(?:[.,]\d+)?)", block, re.IGNORECASE)
    if gentleman_m:
        val = _num(gentleman_m.group(1))
        if val is not None and 0 <= val <= 120:
            data.gentleman = val
    # Здравомыслие — реалистичность, трезвость, практичность суждений
    sanity_m = re.search(r"\b(?:здравомыслие|здравый\s+смысл|common[\s\-]?sense|sanity)\b[ \t]*[:=—\-][ \t]*(\d{1,3}(?:[.,]\d+)?)", block, re.IGNORECASE)
    if sanity_m:
        val = _num(sanity_m.group(1))
        if val is not None and 0 <= val <= 120:
            data.sanity = val

    # Валидность (автоматическая оценка)
    if data.F is not None and data.F >= 100:
        data.validity = "invalid"
    elif data.F is not None and data.F >= 80:
        data.validity = "questionable"
    elif data.L is not None and data.K is not None:
        # Классическая подсказка: L+K низкие + F высокий = профиль сомнителен
        if data.L <= 40 and data.K <= 40 and data.F is not None and data.F >= 65:
            data.validity = "questionable"
        else:
            data.validity = "valid"
    elif data.F is not None and data.F <= 65:
        data.validity = "valid"

    if found == 0:
        notes.append("MMPI/СМИЛ: ничего не распознано. Проверь формат (T-баллы).")
    elif found < 8:
        notes.append(f"MMPI/СМИЛ: распознано {found}/13 шкал. Возможно, формат нестандартный.")
    return data, notes


def parse_disc(block: str) -> tuple[DISC, list[str]]:
    data = DISC()
    notes = []
    # Сначала пробуем полные названия факторов
    axis_syns = {
        "D": [r"доминирование", r"доминантность", r"dominance"],
        "I": [r"влияние", r"influence"],
        "S": [r"стабильность", r"steadiness"],
        "C": [r"соответствие", r"компетентность", r"conscientiousness"],
    }
    found = 0
    for k, syns in axis_syns.items():
        for syn in syns:
            pat = rf"(?:{syn})[ \t]*[:=—\-][ \t]*(\d{{1,3}}(?:[.,]\d+)?)"
            m = re.search(pat, block, re.IGNORECASE)
            if m:
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 100:
                    setattr(data, k, val)
                    found += 1
                    break
    # Если не нашли все через полные названия — пробуем одиночные буквы D/I/S/C
    # но только когда вокруг — НЕ буквы (т.е. это отдельный токен)
    for k in ["D", "I", "S", "C"]:
        if getattr(data, k) is None:
            # Ищем одиночную букву, окружённую не-буквами
            pat = rf"(?<![A-Za-zА-Яа-я]){k}(?![A-Za-zА-Яа-я])[ \t]*[:=—\-][ \t]*(\d{{1,3}}(?:[.,]\d+)?)"
            for m in re.finditer(pat, block):
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 100:
                    setattr(data, k, val)
                    found += 1
                    break
    if found == 0:
        notes.append("DISC: ничего не распознано. Проверь формат (баллы 0-100).")
    elif found < 4:
        notes.append(f"DISC: распознано {found}/4 осей. Возможно, формат нестандартный.")
    return data, notes


def parse_holland(block: str) -> tuple[Holland, list[str]]:
    data = Holland()
    notes = []
    axis_syns = {
        "R": [r"реалистич", r"^R\b", r"realistic", r"\bR\b"],
        "I": [r"исследователь", r"^I\b", r"investigative", r"\bI\b"],
        "A": [r"артистич", r"^A\b", r"artistic", r"\bA\b"],
        "S": [r"социальн", r"^S\b", r"social", r"\bS\b"],
        "E": [r"предприимчив", r"enterprising", r"^E\b", r"\bE\b"],
        "C": [r"конвенциональн", r"conventional", r"^C\b", r"\bC\b"],
    }
    found = 0
    for k, syns in axis_syns.items():
        for syn in syns:
            pat = rf"(?:{syn})[ \t]*[:=—\-][ \t]*(\d{{1,3}}(?:[.,]\d+)?)"
            m = re.search(pat, block, re.IGNORECASE)
            if m:
                val = _num(m.group(1))
                if val is not None and 0 <= val <= 100:
                    setattr(data, k, val)
                    found += 1
                    break
    # Код
    code_m = re.search(r"\bКод[:\s]+([RIASEC]{3})\b", block, re.IGNORECASE)
    if code_m:
        data.code = code_m.group(1).upper()
    if found == 0 and not data.code:
        notes.append("HOLLAND: ничего не распознано. Проверь формат (RIASEC).")
    return data, notes


def parse_mbti(block: str) -> tuple[MBTI, list[str]]:
    data = MBTI()
    notes = []
    # Прямой тип вида INTJ, ESFP и т.д.
    m = re.search(r"\b([EISNTFJP]{4})\b", block)
    if m:
        t = m.group(1).upper()
        if len(set(t)) >= 3:  # защита от мусора типа EEEE
            data.type = t
            data.E_I = t[0]
            data.S_N = t[1]
            data.T_F = t[2]
            data.J_P = t[3]
    if not data.type:
        notes.append("MBTI: не удалось распознать 4-буквенный тип (напр. INTJ).")
    return data, notes


def parse_amthauer(block: str) -> tuple[Amthauer, list[str]]:
    data = Amthauer()
    notes = []
    m = re.search(r"(?:IQ|ай?кью|общий)[ \t:]*=?[ \t]*(\d{2,3})", block, re.IGNORECASE)
    if m:
        data.iq = int(m.group(1))
    if not data.iq:
        notes.append("Амтхауэр: не найден общий IQ.")
    return data, notes


# === Главная функция парсинга ===

def parse_text(full: str, slides_count: int = 0) -> ParsedProfile:
    """Общая логика «текст → ParsedProfile».
    Используется и из PPTX-, и из PDF-парсера (для сканов после OCR).
    """
    employee = parse_employee_info(full)
    blocks = detect_method_blocks(full)

    methods = MethodScores()
    notes = []

    if blocks["cattell_16pf"].strip():
        data, n_notes = parse_cattell(blocks["cattell_16pf"])
        methods.cattell_16pf = data
        notes.extend(n_notes)

    if blocks["big_five"].strip():
        data, n_notes = parse_big_five(blocks["big_five"])
        methods.big_five = data
        notes.extend(n_notes)

    if blocks["mmpi"].strip():
        data, n_notes = parse_mmpi(blocks["mmpi"])
        methods.mmpi = data
        notes.extend(n_notes)

    if blocks["disc"].strip():
        data, n_notes = parse_disc(blocks["disc"])
        methods.disc = data
        notes.extend(n_notes)

    if blocks["holland"].strip():
        data, n_notes = parse_holland(blocks["holland"])
        methods.holland = data
        notes.extend(n_notes)

    if blocks["mbti"].strip():
        data, n_notes = parse_mbti(blocks["mbti"])
        methods.mbti = data
        notes.extend(n_notes)

    if blocks["amthauer"].strip():
        data, n_notes = parse_amthauer(blocks["amthauer"])
        methods.amthauer = data
        notes.extend(n_notes)

    # === ЭФКО-методики (15 тестов + ЭЧВ/ЭЧМ) ===
    efko_data, efko_notes = parse_efko(full)
    if efko_data:
        methods.efko = efko_data
        notes.extend(efko_notes)

    if not any([methods.cattell_16pf, methods.big_five, methods.mmpi,
                methods.disc, methods.holland, methods.mbti, methods.amthauer,
                methods.efko]):
        notes.append("⚠️ Ни одна методика не распознана. Возможно, в презентации нестандартный формат. Смотри docs/presentation_template.md.")

    return ParsedProfile(
        employee=employee,
        methods=methods,
        raw_text=full,
        slides_count=slides_count,
        notes=notes,
    )


def parse_pptx(path: str, source_filename: str | None = None) -> ParsedProfile:
    full, _slides, n = extract_all_text(path)
    profile = parse_text(full, slides_count=n)
    from .pptx_structured import enrich_profile_from_pptx
    return enrich_profile_from_pptx(path, profile, source_filename=source_filename)


# === ЭФКО-методики (15 тестов + ЭЧВ/ЭЧМ + блоки из реальных pptx) ===

EFKO_TEST_PATTERNS = {
    "sluzhebnye_otnosheniya_1": [r"Служебные\s+отношения\s*1"],
    "sluzhebnye_otnosheniya_2": [r"Служебные\s+отношения\s*2"],
    "logicheskoe_myshlenie":    [r"Логическое\s+мышление"],
    "leksika":                  [r"^\s*Лексика\s*[\|]?\s*\d"],
    "zhiznennye_paradigmy":     [r"Жизненные\s+парадигмы"],
    "vizualnye_obrazy":         [r"Визуальные\s+образы"],
    "vospriyatie_otnosheniy":   [r"Восприятие\s+отношений"],
    "obraznoe_myshlenie":       [r"Образное\s+мышление"],
    "socialnye_otnosheniya_1":  [r"Социальные\s+отношения\s*1"],
    "sociokulturnyi_vzglyad_1": [r"Социокультурный\s+взгляд\s*1"],
    "socialnye_orientiry":      [r"Социальные\s+ориентиры"],
    "sociokulturnyi_vzglyad_2": [r"Социокультурный\s+взгляд\s*2"],
    "organizaciya_truda":       [r"Организация\s+труда"],
    "predpochteniya_v_deyatelnosti": [r"Предпочтения\s+в\s+деятельности"],
    "socialnye_otnosheniya_2":  [r"Социальные\s+отношения\s*2"],
}


def _num_optional(s):
    """Достать первое число из строки (в т.ч. '7,5' → 7.5). None если пусто/дефис."""
    s = (s or "").replace(",", ".").strip()
    if not s or s in ("-", "—", ""):
        return None
    m = re.search(r"[-+]?\d+(?:\.\d+)?", s)
    if not m:
        return None
    try:
        return float(m.group(0))
    except ValueError:
        return None


def parse_efko(full: str):
    """Парсер ЭФКО-методик. Возвращает (EFKOSet | None, notes)."""
    from .methods.efko import (
        EFKOSet, SluzhebnyeOtnosheniya1, SluzhebnyeOtnosheniya2,
        LogicheskoeMyshlenie, Leksika, ZhiznennyeParadigmy,
        VizualnyeObrazy, VospriyatieOtnosheniy, ObraznoeMyshlenie,
        SocialnyeOtnosheniya1, SocialnyeOtnosheniya2,
        SociokulturnyiVzglyad1, SociokulturnyiVzglyad2,
        SocialnyeOrientiry, OrganizaciyaTruda,
        PredpochteniyaVDeatelnosti, IntellektEfk, Aktivnost, Empaty, EchvEchm,
    )

    notes = []
    efko = EFKOSet()
    lines = full.splitlines()

    # === ИНТЕЛЛЕКТ (Логический, Образный, Лексика) ===
    intellekt_block = ""
    for i, line in enumerate(lines):
        if re.search(r"^\s*ИНТЕЛЛЕКТ\s*[\|]?\s*%?\s*$", line, re.IGNORECASE):
            intellekt_block = "\n".join(lines[i:min(len(lines), i + 8)])
            break
    intellekt = IntellektEfk()
    if intellekt_block:
        m = re.search(r"Логическ\w*\s*интеллект\s*[:=—\-]?\s*([0-9.,\-]+)", intellekt_block, re.IGNORECASE)
        if m: intellekt.logicheskiy = _num_optional(m.group(1))
        m = re.search(r"Образн\w*\s*интеллект\s*[:=—\-]?\s*([0-9.,\-]+)", intellekt_block, re.IGNORECASE)
        if m: intellekt.obrazny = _num_optional(m.group(1))
        m = re.search(r"^\s*Лексика\s*[:=—\-]?\s*([0-9.,\-]+)", intellekt_block, re.IGNORECASE | re.MULTILINE)
        if m: intellekt.leksika = _num_optional(m.group(1))
    if any(v is not None for v in [intellekt.logicheskiy, intellekt.obrazny, intellekt.leksika]):
        efko.intellekt = intellekt
        notes.append("ЭФКО: распознан блок ИНТЕЛЛЕКТ (Логический, Образный, Лексика).")

    # === АКТИВНОСТЬ (3 шкалы) ===
    aktivnost_block = ""
    for i, line in enumerate(lines):
        if re.search(r"^\s*АКТИВНОСТЬ\s*[\|]?\s*%?\s*$", line, re.IGNORECASE):
            aktivnost_block = "\n".join(lines[i:min(len(lines), i + 8)])
            break
    aktivnost = Aktivnost()
    if aktivnost_block:
        m = re.search(r"Физическ\w*\s*[:=—\-]?\s*([0-9.,\-]+)", aktivnost_block, re.IGNORECASE)
        if m: aktivnost.fizicheskaya = _num_optional(m.group(1))
        m = re.search(r"Интеллектуальн\w*\s*[:=—\-]?\s*([0-9.,\-]+)", aktivnost_block, re.IGNORECASE)
        if m: aktivnost.intellektualnaya = _num_optional(m.group(1))
        m = re.search(r"Коммуникационн\w*\s*[:=—\-]?\s*([0-9.,\-]+)", aktivnost_block, re.IGNORECASE)
        if m: aktivnost.kommunikacionnaya = _num_optional(m.group(1))
    if any(v is not None for v in [aktivnost.fizicheskaya, aktivnost.intellektualnaya, aktivnost.kommunikacionnaya]):
        efko.aktivnost = aktivnost
        notes.append("ЭФКО: распознан блок АКТИВНОСТЬ (Физическая, Интеллектуальная, Коммуникационная).")

    # === ЭМПАТИЯ (сознательная + бессознательная, 2 рода) ===
    empaty_block = ""
    for i, line in enumerate(lines):
        if re.search(r"^\s*ЭМПАТИЯ", line, re.IGNORECASE):
            empaty_block = "\n".join(lines[i:min(len(lines), i + 8)])
            break
    empaty = Empaty()
    if empaty_block:
        # Формат 1: «Сознательная: 55  Бессознательная: 60» (по отдельности)
        m = re.search(r"(?:Сознательная|Созн\.?)\s*[:=—\-]?\s*([0-9.,\-]+)", empaty_block, re.IGNORECASE)
        if m: empaty.conscious = _num_optional(m.group(1))
        m = re.search(r"(?:Бессознательная|Бессозн\.?)\s*[:=—\-]?\s*([0-9.,\-]+)", empaty_block, re.IGNORECASE)
        if m: empaty.unconscious = _num_optional(m.group(1))
        # Формат 2: «Созн. 55 | Бессозн. 60» (через |)
        m = re.search(r"Созн\.\s*([0-9.,]+)\s*[\|]\s*Бессозн\.\s*([0-9.,]+)", empaty_block, re.IGNORECASE)
        if m and empaty.conscious is None:
            empaty.conscious = _num_optional(m.group(1))
            empaty.unconscious = _num_optional(m.group(2))
        # 2 рода
        m2 = re.search(r"2\s*рода\s*,?\s*%?\s*[:=—\-]?\s*([0-9.,]+)?\s*[\|]?\s*([0-9.,]+)?", empaty_block, re.IGNORECASE)
        if m2 and m2.group(1):
            empaty.kind_2_rational = _num_optional(m2.group(1))
            empaty.kind_2_emotional = _num_optional(m2.group(2))
    if any(v is not None for v in [empaty.conscious, empaty.unconscious, empaty.kind_2_rational, empaty.kind_2_emotional]):
        efko.empaty = empaty
        notes.append("ЭФКО: распознан блок ЭМПАТИЯ (сознательная + бессознательная, 2 рода).")

    # === ЭЧВ/ЭЧМ (Глубина) ===
    echv_block = ""
    for i, line in enumerate(lines):
        if re.search(r"ГЛУБИНА\s*ЭЧВ", line, re.IGNORECASE) or re.search(r"Глубина\s+ЭЧВ", line, re.IGNORECASE):
            echv_block = "\n".join(lines[i:min(len(lines), i + 4)])
            break
    echv = EchvEchm()
    if echv_block:
        m = re.search(r"ЭЧВ\s*[:=—\-]?\s*([0-9.,\-]+)", echv_block, re.IGNORECASE)
        if m: echv.echv = _num_optional(m.group(1))
        m = re.search(r"ЭЧМ\s*[:=—\-]?\s*([0-9.,\-]+)", echv_block, re.IGNORECASE)
        if m: echv.echm = _num_optional(m.group(1))
    if echv.echv is not None or echv.echm is not None:
        efko.echv_echm = echv
        notes.append("ЭФКО: распознан блок Глубина ЭЧВ/ЭЧМ.")

    # === Простые однозначные блоки ===
    m = re.search(r"Постмодерн\s*,?\s*%?\s*[:=—\-]?\s*([0-9.,\-]+)", full, re.IGNORECASE)
    if m:
        efko.postmodern = _num_optional(m.group(1))
        notes.append("ЭФКО: распознан блок Постмодерн, %.")

    # Эти поля извлекаются структурным PPTX-парсером. Свободные regex давали
    # ложные захваты соседних заголовков и таблиц.

    # === 15 ЭФКО-тестов (структура для будущих данных) ===
    test_labels_map = {
        "sluzhebnye_otnosheniya_1": (r"Служебные\s+отношения\s*1", SluzhebnyeOtnosheniya1),
        "sluzhebnye_otnosheniya_2": (r"Служебные\s+отношения\s*2", SluzhebnyeOtnosheniya2),
        "logicheskoe_myshlenie":    (r"Логическое\s+мышление", LogicheskoeMyshlenie),
        "leksika":                  (r"^\s*Лексика\s*[\|]?\s*\d", Leksika),
        "zhiznennye_paradigmy":     (r"Жизненные\s+парадигмы", ZhiznennyeParadigmy),
        "vizualnye_obrazy":         (r"Визуальные\s+образы", VizualnyeObrazy),
        "vospriyatie_otnosheniy":   (r"Восприятие\s+отношений", VospriyatieOtnosheniy),
        "obraznoe_myshlenie":       (r"Образное\s+мышление", ObraznoeMyshlenie),
        "socialnye_otnosheniya_1":  (r"Социальные\s+отношения\s*1", SocialnyeOtnosheniya1),
        "sociokulturnyi_vzglyad_1": (r"Социокультурный\s+взгляд\s*1", SociokulturnyiVzglyad1),
        "socialnye_orientiry":      (r"Социальные\s+ориентиры", SocialnyeOrientiry),
        "sociokulturnyi_vzglyad_2": (r"Социокультурный\s+взгляд\s*2", SociokulturnyiVzglyad2),
        "organizaciya_truda":       (r"Организация\s+труда", OrganizaciyaTruda),
        "predpochteniya_v_deyatelnosti": (r"Предпочтения\s+в\s+деятельности", PredpochteniyaVDeatelnosti),
        "socialnye_otnosheniya_2":  (r"Социальные\s+отношения\s*2", SocialnyeOtnosheniya2),
    }
    for key, (pattern, cls) in test_labels_map.items():
        if re.search(pattern, full, re.IGNORECASE):
            setattr(efko, key, cls())

    has_any = any([
        efko.intellekt, efko.aktivnost, efko.empaty, efko.echv_echm,
        efko.postmodern is not None,
        bool(efko.paradigma_k_rabotodatelyu), bool(efko.employer_paradigm),
        bool(efko.life_gamble), bool(efko.harmony_decision),
        bool(efko.work_discomfort), bool(efko.achievement_model),
        bool(efko.personnel_types), bool(efko.safety_attitude),
        any(getattr(efko, k, None) is not None for k in EFKO_TEST_PATTERNS),
    ])
    if not has_any:
        return None, []
    if not notes:
        notes.append("ЭФКО: найдены лейблы 15 тестов, но без баллов (ждём методичку).")
    return efko, notes
